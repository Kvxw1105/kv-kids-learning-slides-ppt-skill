#!/usr/bin/env python3
import argparse
import hashlib
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
LINE_NAME = re.compile(r"(?:line|divider|rule|border|underline|stroke|separator|guide)", re.I)
LABEL_NAME = re.compile(r"(?:label|tag|badge|chip|pill|keyword)", re.I)
SUBJECT_NAME = re.compile(r"(?:flower|sun|cloud|avatar|girl|boy|face|pot|can|drop|hero|character)", re.I)
CONTAINER_NAME = re.compile(r"(?:panel|card|bubble|callout|leaf|ribbon|paper|speech|caption|banner|pill|badge|frame)", re.I)


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def box(shape):
    try:
        return (int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height))
    except Exception:
        return None


def expand(b, pad):
    return (b[0] - pad, b[1] - pad, b[2] + pad, b[3] + pad)


def intersection_area(a, b):
    x1 = max(a[0], b[0]); y1 = max(a[1], b[1])
    x2 = min(a[2], b[2]); y2 = min(a[3], b[3])
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def area(b):
    return max(0, b[2] - b[0]) * max(0, b[3] - b[1])


def is_text(shape):
    return bool(getattr(shape, "has_text_frame", False) and shape.text.strip())


def is_line_like(shape):
    if is_text(shape):
        return False
    name = getattr(shape, "name", "") or ""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return True
    except Exception:
        pass
    b = box(shape)
    if not b:
        return False
    w, h = b[2] - b[0], b[3] - b[1]
    thin = min(w, h) <= int(0.045 * EMU_PER_INCH) and max(w, h) >= int(0.25 * EMU_PER_INCH)
    return bool(LINE_NAME.search(name) or thin)


def inside_slide(b, width, height):
    return b[0] >= 0 and b[1] >= 0 and b[2] <= width and b[3] <= height


def main() -> int:
    p = argparse.ArgumentParser(description="Audit actual PPTX geometry for text-line and label-subject conflicts")
    p.add_argument("pptx")
    p.add_argument("--profile", default="classroom_learning")
    p.add_argument("--clearance-in", type=float, default=0.06)
    p.add_argument("--out")
    args = p.parse_args()

    path = Path(args.pptx)
    errors = []
    warnings = []
    slides = []
    if not path.exists():
        result = {"status": "BLOCKED", "errors": ["file not found"], "warnings": []}
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 1

    prs = Presentation(path)
    pad = int(args.clearance_in * EMU_PER_INCH)
    for slide_index, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        text_shapes = []
        line_shapes = []
        label_shapes = []
        subject_shapes = []
        container_shapes = []
        slide_findings = []
        for shape_index, shape in enumerate(shapes):
            name = getattr(shape, "name", "") or ""
            b = box(shape)
            if not b:
                continue
            if not inside_slide(b, int(prs.slide_width), int(prs.slide_height)) and not re.search(r"(?:blob|background|backdrop|scallop|bleed|edge)", name, re.I):
                errors.append(f"S{slide_index:02d}:{name} extends outside slide bounds")
                slide_findings.append({"severity": "BLOCKED", "type": "out_of_bounds", "shape": name})
            if is_text(shape):
                text_shapes.append((shape_index, name, shape.text.strip()[:120], b))
                if LABEL_NAME.search(name):
                    label_shapes.append((name, b))
            if is_line_like(shape):
                line_shapes.append((name, b))
            if SUBJECT_NAME.search(name) and not re.search(r"(?:backdrop|blob|card|panel|background)", name, re.I):
                subject_shapes.append((name, b))
            if not is_text(shape) and CONTAINER_NAME.search(name) and not re.search(r"(?:background|backdrop|watercolorblob)", name, re.I):
                container_shapes.append((name, b))

        for _, tname, text, tbox in text_shapes:
            safe_box = expand(tbox, pad)
            for lname, lbox in line_shapes:
                if tname == lname:
                    continue
                if intersection_area(safe_box, expand(lbox, max(1, pad // 3))) > 0:
                    msg = f"S{slide_index:02d}:{tname} ('{text[:36]}') is too close to or intersects {lname}"
                    errors.append(msg)
                    slide_findings.append({"severity": "BLOCKED", "type": "text_line_clearance", "text_shape": tname, "line_shape": lname})

        for lname, lbox in label_shapes:
            larea = max(1, area(lbox))
            for sname, sbox in subject_shapes:
                if lname == sname:
                    continue
                ratio = intersection_area(lbox, sbox) / larea
                if ratio > 0.12:
                    msg = f"S{slide_index:02d}:{lname} overlaps {sname} by {ratio:.0%} of label area"
                    errors.append(msg)
                    slide_findings.append({"severity": "BLOCKED", "type": "label_subject_overlap", "label": lname, "subject": sname, "ratio": round(ratio, 3)})

        for tindex, tname, text, tbox in text_shapes:
            tarea = max(1, area(tbox))
            for later in shapes[tindex + 1:]:
                if is_text(later) or is_line_like(later):
                    continue
                lname = getattr(later, "name", "") or ""
                lbox = box(later)
                if not lbox:
                    continue
                ratio = intersection_area(tbox, lbox) / tarea
                if ratio > 0.25:
                    msg = f"S{slide_index:02d}:{tname} ('{text[:36]}') is occluded by later shape {lname} ({ratio:.0%})"
                    errors.append(msg)
                    slide_findings.append({"severity": "BLOCKED", "type": "text_occlusion", "text_shape": tname, "covering_shape": lname, "ratio": round(ratio, 3)})

        inner_pad = int(0.02 * EMU_PER_INCH)
        for _, tname, text, tbox in text_shapes:
            tarea = max(1, area(tbox))
            candidates = []
            for cname, cbox in container_shapes:
                ratio = intersection_area(tbox, cbox) / tarea
                if ratio >= 0.55:
                    candidates.append((ratio, area(cbox), cname, cbox))
            if candidates:
                _, _, cname, cbox = sorted(candidates, key=lambda item: (-item[0], item[1]))[0]
                safe = (cbox[0] + inner_pad, cbox[1] + inner_pad, cbox[2] - inner_pad, cbox[3] - inner_pad)
                if not (tbox[0] >= safe[0] and tbox[1] >= safe[1] and tbox[2] <= safe[2] and tbox[3] <= safe[3]):
                    msg = f"S{slide_index:02d}:{tname} ('{text[:36]}') extends beyond or touches container {cname}"
                    errors.append(msg)
                    slide_findings.append({"severity": "BLOCKED", "type": "text_container_clearance", "text_shape": tname, "container_shape": cname})

        slides.append({"slide": slide_index, "findings": slide_findings, "text_shapes": len(text_shapes), "line_shapes": len(line_shapes), "container_shapes": len(container_shapes)})

    result = {
        "status": "PASS" if not errors else "BLOCKED",
        "profile": args.profile,
        "pptx_sha256": sha256_file(path),
        "clearance_inches": args.clearance_in,
        "errors": errors,
        "warnings": warnings,
        "slides": slides,
    }
    text = json.dumps(result, ensure_ascii=False, indent=2)
    if args.out:
        Path(args.out).write_text(text, encoding="utf-8")
    print(text)
    return 0 if result["status"] == "PASS" else 1


if __name__ == "__main__":
    raise SystemExit(main())
