#!/usr/bin/env python3
import json
import subprocess
import sys
import tempfile
import hashlib
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
SPEECH = ROOT / "assets/sample-lessons/child-speech-flower"
CLASSROOM = ROOT / "assets/sample-lessons/remainder-division"


def run(args: list[str]) -> tuple[int, str, str]:
    completed = subprocess.run(args, capture_output=True, text=True)
    return completed.returncode, completed.stdout, completed.stderr


def check(label: str, args: list[str], expect_zero: bool = True) -> None:
    code, stdout, stderr = run(args)
    if expect_zero and code != 0:
        raise RuntimeError(f"{label} failed ({code}): {stderr or stdout}")
    print(f"PASS {label}")


def main() -> int:
    check("speech spec", [sys.executable, str(SCRIPTS / "validate_slide_spec.py"), str(SPEECH / "slide_spec.json")])
    check("speech typography", [sys.executable, str(SCRIPTS / "validate_typography.py"), str(SPEECH / "slide_spec.json")])
    check("speech visual plan", [sys.executable, str(SCRIPTS / "validate_visual_plan.py"), str(SPEECH / "visual_plan.json")])
    check("classroom spec", [sys.executable, str(SCRIPTS / "validate_slide_spec.py"), str(CLASSROOM / "slide_spec.json")])
    check("classroom visual plan", [sys.executable, str(SCRIPTS / "validate_visual_plan.py"), str(CLASSROOM / "visual_plan.json")])

    code, stdout, stderr = run(
        [
            sys.executable,
            str(SCRIPTS / "query_layouts.py"),
            "--age",
            "4-6",
            "--usage",
            "child_stage_speech",
            "--stage-safe",
        ]
    )
    if code != 0:
        raise RuntimeError(stderr or stdout)
    layouts = json.loads(stdout)
    if len(layouts) < 8:
        raise RuntimeError(f"expected at least 8 stage-safe layouts, found {len(layouts)}")
    print(f"PASS layout registry ({len(layouts)} child speech layouts)")

    with tempfile.TemporaryDirectory() as temp:
        temp_path = Path(temp)
        composed = temp_path / "slide_spec_composed.json"
        planned = temp_path / "slide_spec_planned.json"
        motion = temp_path / "motion_plan.json"
        stickers = temp_path / "stickers"
        storybook_components = temp_path / "storybook_components"
        stage = temp_path / "stage.pptx"
        script = temp_path / "script.md"
        script_docx = temp_path / "script.docx"
        check(
            "visual component composer",
            [sys.executable, str(SCRIPTS / "compose_visual_components.py"), str(SPEECH / "slide_spec.json"), "--out", str(composed)],
        )
        check("responsive storybook planner", [sys.executable, str(SCRIPTS / "plan_storybook_layout.py"), str(composed), "--out", str(planned), "--density", "balanced"])
        check("visual density", [sys.executable, str(SCRIPTS / "validate_visual_density.py"), str(planned)])
        check("storybook composition", [sys.executable, str(SCRIPTS / "validate_storybook_composition.py"), str(planned)])
        check("composed typography", [sys.executable, str(SCRIPTS / "validate_typography.py"), str(planned)])
        check("sticker pack", [sys.executable, str(SCRIPTS / "build_sticker_pack.py"), "--out-dir", str(stickers)])
        check("storybook component pack", [sys.executable, str(SCRIPTS / "build_storybook_components.py"), "--out-dir", str(storybook_components), "--palette", "watercolor"])
        check("motion plan", [sys.executable, str(SCRIPTS / "build_motion_plan.py"), str(planned), "--out", str(motion)])
        check(
            "speech fallback builder",
            [
                sys.executable,
                str(SCRIPTS / "build_child_speech_pptx.py"),
                str(planned),
                "--stage-out",
                str(stage),
                "--script-out",
                str(script),
                "--script-docx-out",
                str(script_docx),
            ],
        )
        check(
            "speaker script Markdown",
            [
                sys.executable,
                str(SCRIPTS / "validate_speaker_script_md.py"),
                str(script),
            ],
        )
        check(
            "speaker script DOCX",
            [
                sys.executable,
                str(SCRIPTS / "inspect_speaker_script_docx.py"),
                str(script_docx),
                "--markdown",
                str(script),
            ],
        )
        check(
            "speech PPTX inspection",
            [
                sys.executable,
                str(SCRIPTS / "inspect_pptx.py"),
                str(stage),
                "--profile",
                "child_stage_speech",
            ],
        )
        check(
            "speech PPTX geometry",
            [sys.executable, str(SCRIPTS / "audit_pptx_geometry.py"), str(stage), "--profile", "child_stage_speech"],
        )
        review_image = temp_path / "review-slide-1.png"
        from PIL import Image
        Image.new("RGB", (64, 36), "white").save(review_image)
        h = hashlib.sha256(stage.read_bytes()).hexdigest()
        review = temp_path / "visual_review.json"
        review.write_text(json.dumps({
            "pptx_sha256": h,
            "slides": [{"slide": 1, "image": str(review_image), "status": "PASS", "issues": [], "notes": "checked"}]
        }), encoding="utf-8")
        check(
            "visual review validator",
            [sys.executable, str(SCRIPTS / "validate_visual_review.py"), str(review), "--pptx", str(stage)],
        )

    with tempfile.TemporaryDirectory() as temp:
        bad = Path(temp) / "bad_density.json"
        bad.write_text(json.dumps({
            "slides": [{
                "id": "S01",
                "screen": {"text_blocks": [{"text": "Hello", "lang": "en", "font_size_pt": 30}]},
                "visual": {"layout": "show-one-object", "completion": {"minimum_layers": 4, "base_layers": ["background", "text"]}}
            }]
        }), encoding="utf-8")
        code, stdout, stderr = run([sys.executable, str(SCRIPTS / "validate_visual_density.py"), str(bad)])
        if code == 0:
            raise RuntimeError("visual density negative test should block a text-only unfinished slide")
        print("PASS visual density negative test")

    with tempfile.TemporaryDirectory() as temp:
        bad = Path(temp) / "bad_typography.json"
        bad.write_text(json.dumps({
            "usage_mode": "child_stage_speech",
            "visual_system": "storybook-watercolor-glow",
            "typography": {"theme_id": "watercolor-storybook"},
            "slides": [{
                "id": "S01",
                "screen": {"text_blocks": [{"id": "zh", "role": "translation", "text": "sample", "lang": "zh", "font_size_pt": 24, "color": "#000000"}]},
                "visual": {"completion": {"density": "balanced", "minimum_layers": 4, "base_layers": ["background", "text"], "white_space_intent": "test", "fill_strategy": ["enlarge type"]}, "native_scene": True}
            }]
        }), encoding="utf-8")
        code, stdout, stderr = run([sys.executable, str(SCRIPTS / "validate_typography.py"), str(bad)])
        if code == 0:
            raise RuntimeError("typography negative test should block pure black child-stage text")
        print("PASS typography negative test")


    with tempfile.TemporaryDirectory() as temp:
        bad = Path(temp) / "bad_storybook.json"
        bad.write_text(json.dumps({
            "usage_mode": "child_stage_speech",
            "slides": [{
                "id": "S01",
                "screen": {"text_blocks": [
                    {"id": "en", "role": "core_line", "text": "This is my flower", "lang": "en", "font_size_pt": 34},
                    {"id": "zh", "role": "translation", "text": "这是我的小花", "lang": "zh", "font_size_pt": 14},
                    {"id": "k1", "role": "keyword", "text": "pink", "lang": "en", "font_size_pt": 20},
                    {"id": "k2", "role": "keyword", "text": "small", "lang": "en", "font_size_pt": 20},
                    {"id": "k3", "role": "keyword", "text": "soft", "lang": "en", "font_size_pt": 20}
                ]},
                "visual": {"layout_plan": {"archetype": "picturebook-window", "content_level": "sparse", "hero_scale_ratio": 0.2, "text_container_style": "rectangle", "zones": {"hero": [0,0,4,4]}, "keyword_group": {"style": "none", "block_ids": []}}}
            }]
        }), encoding="utf-8")
        code, stdout, stderr = run([sys.executable, str(SCRIPTS / "validate_storybook_composition.py"), str(bad)])
        if code == 0:
            raise RuntimeError("storybook composition negative test should block orphan labels and weak bilingual hierarchy")
        print("PASS storybook composition negative test")


    with tempfile.TemporaryDirectory() as temp:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        bad_pptx = Path(temp) / "bad_geometry.pptx"
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.7))
        tb.name = "S01_Title_EN"
        tb.text_frame.paragraphs[0].text = "Sunshine and gentle rain"
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
        line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.45), Inches(6.4), Inches(0.02))
        line.name = "S01_Title_Divider_Line"
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(120, 150, 120)
        line.line.fill.background()
        prs.save(bad_pptx)
        code, stdout, stderr = run([sys.executable, str(SCRIPTS / "audit_pptx_geometry.py"), str(bad_pptx), "--profile", "child_stage_speech"])
        if code == 0:
            raise RuntimeError("geometry negative test should block a text-line intersection")
        print("PASS geometry negative test")

    with tempfile.TemporaryDirectory() as temp:
        missing = Path(temp) / "visual_review.json"
        missing.write_text(json.dumps({"pptx_sha256": "wrong", "slides": [{"slide": 1, "image": "missing.png", "status": "PENDING"}]}), encoding="utf-8")
        dummy = Path(temp) / "dummy.pptx"
        dummy.write_bytes(b"not-a-real-pptx")
        code, stdout, stderr = run([sys.executable, str(SCRIPTS / "validate_visual_review.py"), str(missing), "--pptx", str(dummy)])
        if code == 0:
            raise RuntimeError("visual review negative test should block stale or pending review")
        print("PASS visual review negative test")

    print("ALL TESTS PASSED")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
