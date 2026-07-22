#!/usr/bin/env python3
import argparse
import hashlib
import json
import re
import sys
import zipfile
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PATH_PATTERNS = [
    re.compile(r"/mnt/data/[^\s\"'<>]+"),
    re.compile(r"/home/[^\s\"'<>]+"),
    re.compile(r"[A-Za-z]:\\Users\\[^\s\"'<>]+"),
    re.compile(r"file:///[^\s\"'<>]+"),
]
GENERIC_NAME = re.compile(
    r"^(?:TextBox|Text|Picture|Image|Rectangle|Rounded Rectangle|Oval|Line|Shape|Group|Freeform|Object)(?:\s|_|-)*\d*$",
    re.IGNORECASE,
)


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def extract_font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(round(run.font.size.pt, 2))
    return sizes




def shape_box(shape):
    try:
        return (int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height))
    except Exception:
        return None


def intersection_area(a, b) -> int:
    x1 = max(a[0], b[0]); y1 = max(a[1], b[1])
    x2 = min(a[2], b[2]); y2 = min(a[3], b[3])
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def box_area(a) -> int:
    return max(0, a[2] - a[0]) * max(0, a[3] - a[1])



def coverage_metrics(shapes, slide_width: int, slide_height: int, cols: int = 40, rows: int = 23) -> tuple[float, int, float]:
    """Approximate occupied area and largest connected empty region."""
    occupied = [[False for _ in range(cols)] for _ in range(rows)]
    meaningful = 0
    slide_area = max(1, slide_width * slide_height)
    for shape in shapes:
        name = (getattr(shape, "name", "") or "").lower()
        if any(token in name for token in ("blob", "background")):
            continue
        box = shape_box(shape)
        if not box:
            continue
        area = box_area(box)
        if area / slide_area > 0.88:
            continue
        meaningful += 1
        x1, y1, x2, y2 = box
        c1 = max(0, min(cols - 1, int(x1 / slide_width * cols)))
        c2 = max(0, min(cols - 1, int(max(x1, x2 - 1) / slide_width * cols)))
        r1 = max(0, min(rows - 1, int(y1 / slide_height * rows)))
        r2 = max(0, min(rows - 1, int(max(y1, y2 - 1) / slide_height * rows)))
        for row in range(r1, r2 + 1):
            for col in range(c1, c2 + 1):
                cx = (col + 0.5) / cols * slide_width
                cy = (row + 0.5) / rows * slide_height
                if x1 <= cx <= x2 and y1 <= cy <= y2:
                    occupied[row][col] = True
    cells = sum(1 for row in occupied for value in row if value)
    visited = [[False for _ in range(cols)] for _ in range(rows)]
    largest_empty = 0
    for row in range(rows):
        for col in range(cols):
            if occupied[row][col] or visited[row][col]:
                continue
            stack = [(row, col)]
            visited[row][col] = True
            region = 0
            while stack:
                current_row, current_col = stack.pop()
                region += 1
                for delta_row, delta_col in ((1, 0), (-1, 0), (0, 1), (0, -1)):
                    next_row = current_row + delta_row
                    next_col = current_col + delta_col
                    if 0 <= next_row < rows and 0 <= next_col < cols and not occupied[next_row][next_col] and not visited[next_row][next_col]:
                        visited[next_row][next_col] = True
                        stack.append((next_row, next_col))
            largest_empty = max(largest_empty, region)
    total = cols * rows
    return cells / total, meaningful, largest_empty / total

def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect PPTX structure and release risks")
    parser.add_argument("pptx")
    parser.add_argument("--profile", default="classroom_learning")
    parser.add_argument("--max-size-mb", type=float)
    parser.add_argument("--min-font-pt", type=float)
    parser.add_argument("--min-coverage", type=float)
    parser.add_argument("--warn-coverage", type=float)
    parser.add_argument("--max-empty-region", type=float)
    args = parser.parse_args()

    path = Path(args.pptx)
    errors: list[str] = []
    warnings: list[str] = []
    if not path.exists():
        result = {"status": "BLOCKED", "errors": ["file not found"], "warnings": []}
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 1

    max_size_mb = args.max_size_mb
    if max_size_mb is None:
        max_size_mb = 5.0 if args.profile == "child_stage_speech" else 25.0
    min_font_pt = args.min_font_pt
    if min_font_pt is None and args.profile == "child_stage_speech":
        min_font_pt = 20.0
    min_coverage = args.min_coverage if args.min_coverage is not None else (0.28 if args.profile == "child_stage_speech" else 0.0)
    warn_coverage = args.warn_coverage if args.warn_coverage is not None else (0.40 if args.profile == "child_stage_speech" else 0.0)
    max_empty_region = args.max_empty_region if args.max_empty_region is not None else (0.44 if args.profile == "child_stage_speech" else 1.0)

    file_size = path.stat().st_size
    file_size_mb = file_size / (1024 * 1024)
    if file_size_mb > max_size_mb:
        warnings.append(f"file size {file_size_mb:.2f}MB exceeds budget {max_size_mb:.2f}MB")
    if file_size_mb > max_size_mb * 2:
        errors.append(f"file size {file_size_mb:.2f}MB is more than twice the release budget")

    media: list[dict] = []
    duplicate_groups: list[dict] = []
    leaks: list[str] = []
    transition_slides: set[str] = set()
    animation_slides: set[str] = set()
    try:
        with zipfile.ZipFile(path) as archive:
            bad_member = archive.testzip()
            if bad_member:
                errors.append(f"corrupt member: {bad_member}")
            media_hashes: dict[str, list[tuple[str, int]]] = defaultdict(list)
            for info in archive.infolist():
                payload = archive.read(info.filename)
                if info.filename.startswith("ppt/media/"):
                    digest = sha256_bytes(payload)
                    media_hashes[digest].append((info.filename, len(payload)))
                    media.append({"file": info.filename, "bytes": len(payload), "sha256": digest})
                if info.filename.endswith((".xml", ".rels")):
                    text = payload.decode("utf-8", errors="ignore")
                    for pattern in PATH_PATTERNS:
                        for match in pattern.findall(text):
                            leaks.append(f"{info.filename}: {match[:180]}")
                    if info.filename.startswith("ppt/slides/slide") and "<p:transition" in text:
                        transition_slides.add(info.filename)
                    if info.filename.startswith("ppt/slides/slide") and "<p:timing" in text:
                        animation_slides.add(info.filename)
            for digest, items in media_hashes.items():
                if len(items) > 1:
                    duplicate_groups.append(
                        {
                            "sha256": digest,
                            "files": [name for name, _ in items],
                            "duplicate_bytes": sum(size for _, size in items[1:]),
                        }
                    )
    except Exception as exc:
        errors.append(f"zip inspection failed: {exc}")

    if leaks:
        errors.append(f"local path metadata leak found in {len(leaks)} XML location(s)")
    duplicate_bytes = sum(group["duplicate_bytes"] for group in duplicate_groups)
    if duplicate_groups:
        warnings.append(
            f"{len(duplicate_groups)} duplicate media hash group(s), {duplicate_bytes / 1024:.1f}KB avoidable bytes"
        )

    slides_report: list[dict] = []
    generic_names: list[str] = []
    semantic_names = 0
    all_font_sizes: list[float] = []
    notes_count = 0
    editable_text_shapes = 0
    image_count = 0
    label_overlaps: list[str] = []
    try:
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides, 1):
            names: list[str] = []
            texts: list[str] = []
            slide_font_sizes: list[float] = []
            slide_images = 0
            slide_text_shapes = 0
            rectangular_shapes = 0
            playful_shapes = 0
            slide_shapes = list(slide.shapes)
            density_ratio, meaningful_shape_count, largest_empty_ratio = coverage_metrics(
                slide_shapes, int(presentation.slide_width), int(presentation.slide_height)
            )
            for shape in slide_shapes:
                name = getattr(shape, "name", "") or ""
                names.append(name)
                try:
                    auto_name = str(getattr(shape, "auto_shape_type", "") or "").lower()
                except Exception:
                    auto_name = ""
                if any(token in auto_name for token in ("rect", "square")):
                    rectangular_shapes += 1
                if any(token in auto_name for token in ("cloud", "callout", "leaf", "flower", "heart", "star", "sun", "wave", "hexagon", "oval")):
                    playful_shapes += 1
                if GENERIC_NAME.match(name.strip()):
                    generic_names.append(f"S{slide_index:02d}:{name}")
                else:
                    semantic_names += 1
                if getattr(shape, "has_text_frame", False):
                    slide_text_shapes += 1
                    editable_text_shapes += 1
                    if shape.text.strip():
                        texts.append(shape.text.strip()[:100])
                    sizes = extract_font_sizes(shape)
                    slide_font_sizes.extend(sizes)
                    all_font_sizes.extend(sizes)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_images += 1
                    image_count += 1
            if args.profile == "child_stage_speech":
                label_shapes = []
                subject_shapes = []
                for candidate in slide_shapes:
                    cname = getattr(candidate, "name", "") or ""
                    cbox = shape_box(candidate)
                    if not cbox:
                        continue
                    if cname.endswith("_Label") or "_Label_" in cname:
                        label_shapes.append((cname, cbox))
                    elif any(token in cname for token in ("Flower", "Sun", "Cloud", "Avatar", "Pot", "Can", "Drop")):
                        if "Backdrop" not in cname and "Blob" not in cname and "Card" not in cname:
                            subject_shapes.append((cname, cbox))
                for lname, lbox in label_shapes:
                    larea = box_area(lbox) or 1
                    for sname, sbox in subject_shapes:
                        if lname == sname:
                            continue
                        ratio = intersection_area(lbox, sbox) / larea
                        if ratio > 0.12:
                            label_overlaps.append(f"S{slide_index:02d}:{lname} overlaps {sname} by {ratio:.0%} of label area")

            notes = ""
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
            if notes:
                notes_count += 1
            slides_report.append(
                {
                    "slide": slide_index,
                    "shapes": len(slide.shapes),
                    "editable_text_shapes": slide_text_shapes,
                    "images": slide_images,
                    "min_explicit_font_pt": min(slide_font_sizes) if slide_font_sizes else None,
                    "has_notes": bool(notes),
                    "content_coverage_ratio": round(density_ratio, 3),
                    "meaningful_shape_count": meaningful_shape_count,
                    "rectangular_shape_count": rectangular_shapes,
                    "playful_shape_count": playful_shapes,
                    "largest_empty_region_ratio": round(largest_empty_ratio, 3),
                    "text_samples": texts[:3],
                    "shape_names": names[:8],
                }
            )
        slide_count = len(presentation.slides)
    except Exception as exc:
        errors.append(f"PowerPoint parse failed: {exc}")
        slide_count = 0

    total_named = len(generic_names) + semantic_names
    if total_named and len(generic_names) / total_named > 0.5:
        warnings.append(f"{len(generic_names)}/{total_named} shapes use generic names")
    if min_font_pt is not None and all_font_sizes:
        too_small = sorted(size for size in all_font_sizes if size < min_font_pt)
        if too_small:
            errors.append(
                f"{len(too_small)} explicit text run(s) are below {min_font_pt:g}pt; minimum is {min(too_small):g}pt"
            )
    elif min_font_pt is not None and not all_font_sizes:
        warnings.append("no explicit font sizes found; stage readability could not be verified")

    if args.profile == "child_stage_speech":
        for slide_info in slides_report:
            ratio = slide_info.get("content_coverage_ratio", 0)
            count = slide_info.get("meaningful_shape_count", 0)
            empty_ratio = slide_info.get("largest_empty_region_ratio", 1.0)
            if ratio < min_coverage:
                if count < 7:
                    errors.append(f"slide {slide_info['slide']} looks unfinished: coverage {ratio:.0%}, meaningful shapes {count}")
                else:
                    warnings.append(f"slide {slide_info['slide']} has low occupied coverage ({ratio:.0%}) despite multiple small objects")
            elif ratio < warn_coverage and count < 9:
                warnings.append(f"slide {slide_info['slide']} may be visually sparse: coverage {ratio:.0%}, meaningful shapes {count}")
            rects = slide_info.get("rectangular_shape_count", 0)
            playful = slide_info.get("playful_shape_count", 0)
            if rects >= 5 and rects > playful * 2:
                warnings.append(f"slide {slide_info['slide']} is rectangle-dominant ({rects} rectangular vs {playful} playful shapes)")
            if empty_ratio > max_empty_region + 0.15 and ratio < min_coverage and count < 12:
                errors.append(f"slide {slide_info['slide']} has a continuous empty region of {empty_ratio:.0%}")
            elif empty_ratio > max_empty_region and ratio < warn_coverage:
                warnings.append(f"slide {slide_info['slide']} has a large continuous empty region ({empty_ratio:.0%})")
    if args.profile == "child_stage_speech" and notes_count == 0:
        warnings.append("no speaker notes found in PPTX; provide a separate speaker script at minimum")
    if args.profile == "child_stage_speech" and not transition_slides and not animation_slides:
        warnings.append("no transitions or animation timing found; static delivery is valid but reveal plan is unimplemented")
    if label_overlaps:
        errors.append(f"{len(label_overlaps)} label/object overlap(s) found")

    result = {
        "status": "PASS" if not errors else "BLOCKED",
        "profile": args.profile,
        "pptx_sha256": sha256_file(path),
        "file_size_bytes": file_size,
        "file_size_mb": round(file_size_mb, 3),
        "size_budget_mb": max_size_mb,
        "min_coverage_ratio": min_coverage,
        "warn_coverage_ratio": warn_coverage,
        "max_empty_region_ratio": max_empty_region,
        "slide_count": slide_count,
        "editable_text_shapes": editable_text_shapes,
        "image_shapes": image_count,
        "notes_slides": notes_count,
        "min_explicit_font_pt": min(all_font_sizes) if all_font_sizes else None,
        "transition_slide_count": len(transition_slides),
        "animation_slide_count": len(animation_slides),
        "media_files": len(media),
        "duplicate_media_groups": duplicate_groups,
        "local_path_leaks": leaks[:20],
        "generic_shape_names": generic_names[:30],
        "label_overlaps": label_overlaps[:30],
        "errors": errors,
        "warnings": warnings,
        "slides": slides_report,
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["status"] == "PASS" else 1


if __name__ == "__main__":
    raise SystemExit(main())
