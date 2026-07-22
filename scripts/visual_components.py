#!/usr/bin/env python3
"""Native PowerPoint visual-component renderer.

The component plan is intentionally small and deterministic. It generates
editable PowerPoint shapes for child-friendly stickers, ribbons, borders, and
visual cues so a deck does not depend on raster AI illustrations for every
visual detail.
"""
from __future__ import annotations

import math
from pathlib import Path
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_COLORS = {
    "pink": "F6A6C8",
    "deep_pink": "D95D8A",
    "yellow": "FFD66B",
    "green": "88B96B",
    "deep_green": "4E7B55",
    "blue": "9ED9F6",
    "deep_blue": "4D84A8",
    "cream": "FFF9EC",
    "white": "FFFFFF",
    "ink": "503F45",
    "lavender": "C9B8F4",
    "orange": "F3A861",
}


def rgb(value: str) -> RGBColor:
    h = value.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_meta(shape, name: str, descr: str = "") -> None:
    shape.name = name
    for attr in ("nvSpPr", "nvPicPr"):
        node = getattr(shape._element, attr, None)
        if node is not None:
            c_nv_pr = node.cNvPr
            if descr:
                c_nv_pr.set("descr", descr)
            return


def add_shape(slide, shape_type, x, y, w, h, fill, name, line=None, line_width=1.0, rotation=0, transparency=0):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = int(max(0, min(100, transparency)))
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    if rotation:
        shape.rotation = rotation
    set_meta(shape, name)
    return shape


def add_text(slide, text, x, y, w, h, size, name, color="503F45", bold=True, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_meta(box, name)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial Rounded MT Bold"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def _color(item: dict[str, Any], key: str, default: str) -> str:
    value = item.get(key)
    if value:
        return value.lstrip("#")
    return DEFAULT_COLORS.get(default, default)


def draw_star(slide, item, sid):
    return add_shape(slide, MSO_SHAPE.STAR_5_POINT, item["x"], item["y"], item["w"], item["h"], _color(item, "fill", "yellow"), f"{sid}_{item['id']}_Star", line=_color(item, "line", "orange"), line_width=0.8, rotation=item.get("rotation", 0), transparency=item.get("transparency", 0))


def draw_heart(slide, item, sid):
    return add_shape(slide, MSO_SHAPE.HEART, item["x"], item["y"], item["w"], item["h"], _color(item, "fill", "pink"), f"{sid}_{item['id']}_Heart", line=_color(item, "line", "deep_pink"), line_width=0.8, rotation=item.get("rotation", 0), transparency=item.get("transparency", 0))


def draw_leaf_sprig(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    stem = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + w * 0.15), Inches(y + h * 0.85), Inches(x + w * 0.85), Inches(y + h * 0.15))
    stem.line.color.rgb = rgb(_color(item, "line", "deep_green"))
    stem.line.width = Pt(max(1.5, min(w, h) * 3.2))
    set_meta(stem, f"{sid}_{item['id']}_Stem")
    for i, (px, py, rot) in enumerate(((0.20, 0.58, -35), (0.38, 0.43, 28), (0.53, 0.28, -35), (0.68, 0.12, 28)), 1):
        add_shape(slide, MSO_SHAPE.OVAL, x + w * px, y + h * py, w * 0.28, h * 0.22, _color(item, "fill", "green"), f"{sid}_{item['id']}_Leaf_{i}", rotation=rot, transparency=item.get("transparency", 0))


def draw_flower_sticker(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    petal = _color(item, "fill", "pink")
    center = _color(item, "accent", "yellow")
    for i, angle in enumerate(range(0, 360, 60), 1):
        rad = math.radians(angle)
        px = x + w * 0.5 + math.cos(rad) * w * 0.22 - w * 0.16
        py = y + h * 0.5 + math.sin(rad) * h * 0.22 - h * 0.16
        add_shape(slide, MSO_SHAPE.OVAL, px, py, w * 0.32, h * 0.32, petal, f"{sid}_{item['id']}_Petal_{i}", transparency=item.get("transparency", 0))
    add_shape(slide, MSO_SHAPE.OVAL, x + w * 0.35, y + h * 0.35, w * 0.30, h * 0.30, center, f"{sid}_{item['id']}_Center", transparency=item.get("transparency", 0))


def draw_butterfly(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    fill = _color(item, "fill", "lavender")
    accent = _color(item, "accent", "pink")
    add_shape(slide, MSO_SHAPE.OVAL, x, y + h * 0.05, w * 0.42, h * 0.46, fill, f"{sid}_{item['id']}_Wing_L1", rotation=-22)
    add_shape(slide, MSO_SHAPE.OVAL, x + w * 0.05, y + h * 0.45, w * 0.34, h * 0.38, accent, f"{sid}_{item['id']}_Wing_L2", rotation=18)
    add_shape(slide, MSO_SHAPE.OVAL, x + w * 0.58, y + h * 0.05, w * 0.42, h * 0.46, fill, f"{sid}_{item['id']}_Wing_R1", rotation=22)
    add_shape(slide, MSO_SHAPE.OVAL, x + w * 0.61, y + h * 0.45, w * 0.34, h * 0.38, accent, f"{sid}_{item['id']}_Wing_R2", rotation=-18)
    add_shape(slide, MSO_SHAPE.OVAL, x + w * 0.45, y + h * 0.18, w * 0.10, h * 0.60, _color(item, "line", "ink"), f"{sid}_{item['id']}_Body")


def draw_sparkle_cluster(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    coords = [(0.05, 0.32, 0.28), (0.38, 0.02, 0.20), (0.62, 0.42, 0.32), (0.30, 0.65, 0.16)]
    for i, (px, py, scale) in enumerate(coords, 1):
        add_shape(slide, MSO_SHAPE.STAR_5_POINT, x + w * px, y + h * py, w * scale, h * scale, _color(item, "fill", "yellow"), f"{sid}_{item['id']}_Sparkle_{i}", rotation=18 * i, transparency=item.get("transparency", 0))


def draw_ribbon(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    fill = _color(item, "fill", "pink")
    accent = _color(item, "accent", "deep_pink")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + w * 0.08, y, w * 0.84, h, fill, f"{sid}_{item['id']}_Ribbon_Body", line=accent, line_width=0.8)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y + h * 0.12, w * 0.18, h * 0.76, accent, f"{sid}_{item['id']}_Ribbon_Left", rotation=270)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + w * 0.82, y + h * 0.12, w * 0.18, h * 0.76, accent, f"{sid}_{item['id']}_Ribbon_Right", rotation=90)
    if item.get("text"):
        add_text(slide, str(item["text"]), x + w * 0.12, y + h * 0.04, w * 0.76, h * 0.88, float(item.get("font_size_pt", 18)), f"{sid}_{item['id']}_Ribbon_Text", color=_color(item, "text_color", "white"))


def draw_scallop_strip(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    fill = _color(item, "fill", "pink")
    count = max(4, int(item.get("count", 12)))
    diameter = w / count
    for i in range(count):
        add_shape(slide, MSO_SHAPE.OVAL, x + i * diameter, y, diameter * 1.08, h, fill, f"{sid}_{item['id']}_Scallop_{i+1}", transparency=item.get("transparency", 0))


def draw_dots(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    colors = item.get("colors") or [DEFAULT_COLORS["pink"], DEFAULT_COLORS["yellow"], DEFAULT_COLORS["blue"], DEFAULT_COLORS["green"]]
    coords = [(0.05, 0.20, 0.15), (0.38, 0.02, 0.11), (0.68, 0.30, 0.18), (0.25, 0.62, 0.12), (0.76, 0.72, 0.10)]
    for i, (px, py, scale) in enumerate(coords, 1):
        add_shape(slide, MSO_SHAPE.OVAL, x + w * px, y + h * py, w * scale, w * scale, str(colors[(i - 1) % len(colors)]).lstrip("#"), f"{sid}_{item['id']}_Dot_{i}", transparency=item.get("transparency", 0))


def draw_grass(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    fill = _color(item, "fill", "green")
    add_shape(slide, MSO_SHAPE.ARC, x, y + h * 0.25, w, h * 0.75, fill, f"{sid}_{item['id']}_Grass_Base", line=fill, line_width=2)
    for i in range(8):
        px = x + w * (i + 0.5) / 8
        blade = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(y + h), Inches(px + (0.08 if i % 2 else -0.08)), Inches(y + h * (0.30 + 0.08 * (i % 3))))
        blade.line.color.rgb = rgb(fill)
        blade.line.width = Pt(2.2)
        set_meta(blade, f"{sid}_{item['id']}_Blade_{i+1}")


def draw_badge(slide, item, sid):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    fill = _color(item, "fill", "yellow")
    accent = _color(item, "accent", "orange")
    add_shape(slide, MSO_SHAPE.STAR_12_POINT, x, y, w, h, fill, f"{sid}_{item['id']}_Badge", line=accent, line_width=1.0)
    if item.get("text"):
        add_text(slide, str(item["text"]), x + w * 0.16, y + h * 0.17, w * 0.68, h * 0.66, float(item.get("font_size_pt", 16)), f"{sid}_{item['id']}_Badge_Text", color=_color(item, "text_color", "ink"))


def draw_paper_panel(slide, item, sid):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        item["x"], item["y"], item["w"], item["h"],
        _color(item, "fill", "white"),
        f"{sid}_{item['id']}_PaperPanel",
        line=_color(item, "line", "pink"),
        line_width=float(item.get("line_width", 1.0)),
        transparency=item.get("transparency", 0),
    )


def draw_animated_gif(slide, item, sid):
    file_value = item.get("file")
    if not file_value:
        return
    path = Path(file_value)
    if not path.exists():
        return
    picture = slide.shapes.add_picture(
        str(path),
        Inches(item["x"]),
        Inches(item["y"]),
        Inches(item["w"]),
        Inches(item["h"]),
    )
    set_meta(picture, f"{sid}_{item['id']}_AnimatedSticker", "Code-generated animated sticker")


DRAWERS = {
    "star": draw_star,
    "heart": draw_heart,
    "leaf_sprig": draw_leaf_sprig,
    "flower_sticker": draw_flower_sticker,
    "butterfly": draw_butterfly,
    "sparkle_cluster": draw_sparkle_cluster,
    "ribbon": draw_ribbon,
    "scallop_strip": draw_scallop_strip,
    "dots": draw_dots,
    "grass": draw_grass,
    "badge": draw_badge,
    "paper_panel": draw_paper_panel,
    "animated_gif": draw_animated_gif,
}


def add_visual_components(slide, slide_spec: dict[str, Any], sid: str, layer: str = "foreground") -> int:
    count = 0
    plan = slide_spec.get("visual", {}).get("component_plan", []) or []
    for index, raw in enumerate(plan, 1):
        item = dict(raw)
        if item.get("layer", "foreground") != layer:
            continue
        item.setdefault("id", f"component_{index:02d}")
        kind = item.get("type")
        drawer = DRAWERS.get(kind)
        if drawer is None:
            continue
        required = ["x", "y", "w", "h"]
        if not all(key in item for key in required):
            continue
        drawer(slide, item, sid)
        count += 1
    return count
