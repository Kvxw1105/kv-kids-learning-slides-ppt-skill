#!/usr/bin/env python3
"""Portable fallback builder for child stage-speech PPTX files.

Creates a clean audience-facing stage deck and, optionally, a rehearsal deck.
The slide specification remains the source of truth. Text, native diagrams, and
speaker notes stay editable. Complex illustrations can be supplied through an
asset manifest.
"""
from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from visual_components import add_visual_components

SLIDE_W = 13.333
SLIDE_H = 7.5
TYPOGRAPHY_REGISTRY = Path(__file__).resolve().parents[1] / 'assets' / 'registries' / 'typography-themes.json'

PALETTE = {
    "cream": "FFF9EC",
    "pink": "F6A6C8",
    "deep_pink": "D95D8A",
    "yellow": "FFD66B",
    "green": "88B96B",
    "deep_green": "4E7B55",
    "blue": "9ED9F6",
    "deep_blue": "4D84A8",
    "brown": "8A5A3B",
    "ink": "503F45",
    "white": "FFFFFF",
    "soft": "FFFDF8",
}


def rgb(hex_value: str) -> RGBColor:
    h = hex_value.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_shape_meta(shape, name: str, descr: str = "") -> None:
    shape.name = name
    try:
        c_nv_pr = shape._element.nvSpPr.cNvPr
    except AttributeError:
        try:
            c_nv_pr = shape._element.nvPicPr.cNvPr
        except AttributeError:
            return
    if descr:
        c_nv_pr.set("descr", descr)
    elif "descr" in c_nv_pr.attrib:
        del c_nv_pr.attrib["descr"]


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    name: str,
    line: str | None = None,
    line_width: float = 1.0,
    rotation: float | None = None,
):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    if rotation is not None:
        shape.rotation = rotation
    set_shape_meta(shape, name)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    name: str,
    lang: str = "en",
    color: str = PALETTE["ink"],
    bold: bool = True,
    align=PP_ALIGN.CENTER,
    font_name: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_meta(box, name)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = font_name or ("Microsoft YaHei" if lang == "zh" else "Arial Rounded MT Bold")
    return box


def load_typography_theme(spec: dict) -> dict:
    themes = json.loads(TYPOGRAPHY_REGISTRY.read_text(encoding="utf-8"))
    explicit = spec.get("typography", {}).get("theme_id")
    visual_system = spec.get("visual_system")
    for theme in themes:
        if explicit and theme.get("id") == explicit:
            return theme
    for theme in themes:
        if visual_system in theme.get("visual_systems", []):
            return theme
    return themes[0]


def typography_style(theme: dict, role: str, lang: str) -> dict:
    if role == "title":
        key = "title_zh" if lang == "zh" else "title_en"
    elif role in {"translation", "core_line"}:
        key = "core_zh" if lang == "zh" else "core_en"
    else:
        key = "label"
    return dict(theme.get(key, {}))


def add_styled_text(
    slide,
    block: dict,
    x: float,
    y: float,
    w: float,
    h: float,
    name: str,
    theme: dict,
    default_size: float,
    align=PP_ALIGN.CENTER,
):
    lang = block.get("lang", "en")
    role = block.get("role", "core_line")
    style = typography_style(theme, role, lang)
    size = float(block.get("font_size_pt") or style.get("size_pt") or default_size)
    color = str(block.get("color") or style.get("color") or PALETTE["ink"]).lstrip("#")
    families = style.get("font_families") or []
    font_name = block.get("font_family") or (families[0] if families else None)
    bold = bool(block.get("bold", style.get("bold", True)))
    effect = block.get("effect") or style.get("effect", "none")
    if effect in {"soft-shadow", "watercolor-echo"}:
        offset = float(style.get("shadow_offset_in", 0.03))
        shadow_color = str(style.get("shadow_color", "F1DCE4")).lstrip("#")
        add_text(
            slide,
            block.get("text", ""),
            x + offset,
            y + offset,
            w,
            h,
            size,
            name + "_SoftEcho",
            lang=lang,
            color=shadow_color,
            bold=bold,
            align=align,
            font_name=font_name,
        )
    return add_text(
        slide,
        block.get("text", ""),
        x,
        y,
        w,
        h,
        size,
        name,
        lang=lang,
        color=color,
        bold=bold,
        align=align,
        font_name=font_name,
    )


def add_soft_background(slide, sid: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(PALETTE["cream"])
    add_shape(slide, MSO_SHAPE.OVAL, -0.55, -0.7, 3.0, 2.0, "FFE2EE", f"{sid}_Blob_TopLeft")
    add_shape(slide, MSO_SHAPE.OVAL, 10.7, 5.85, 3.2, 2.0, "DFF3E1", f"{sid}_Blob_BottomRight")
    add_shape(slide, MSO_SHAPE.OVAL, 11.7, -0.35, 1.7, 1.0, "FFF0B5", f"{sid}_Blob_TopRight")


def add_flower(slide, x: float, y: float, scale: float, sid: str, prefix: str, bloom: bool = True) -> None:
    stem = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + 0.45 * scale),
        Inches(y + 0.78 * scale),
        Inches(x + 0.45 * scale),
        Inches(y + 1.72 * scale),
    )
    stem.line.color.rgb = rgb(PALETTE["deep_green"])
    stem.line.width = Pt(5 * scale)
    set_shape_meta(stem, f"{sid}_{prefix}_Stem")
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.05 * scale, y + 1.08 * scale, 0.55 * scale, 0.32 * scale, PALETTE["green"], f"{sid}_{prefix}_LeafLeft", rotation=-25)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.38 * scale, y + 1.28 * scale, 0.55 * scale, 0.32 * scale, PALETTE["green"], f"{sid}_{prefix}_LeafRight", rotation=25)
    if bloom:
        for angle in range(0, 360, 60):
            rad = math.radians(angle)
            px = x + 0.45 * scale + math.cos(rad) * 0.30 * scale - 0.18 * scale
            py = y + 0.45 * scale + math.sin(rad) * 0.30 * scale - 0.18 * scale
            add_shape(slide, MSO_SHAPE.OVAL, px, py, 0.36 * scale, 0.36 * scale, PALETTE["pink"], f"{sid}_{prefix}_Petal_{angle}")
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.27 * scale, y + 0.27 * scale, 0.36 * scale, 0.36 * scale, PALETTE["yellow"], f"{sid}_{prefix}_Center")
    else:
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.25 * scale, y + 0.23 * scale, 0.40 * scale, 0.60 * scale, PALETTE["pink"], f"{sid}_{prefix}_Bud")
    add_shape(slide, MSO_SHAPE.CAN, x - 0.05 * scale, y + 1.62 * scale, 1.0 * scale, 0.58 * scale, "C7794C", f"{sid}_{prefix}_Pot", line="A85B38")


def add_watering_can(slide, x: float, y: float, scale: float, sid: str) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.25 * scale, 0.9 * scale, PALETTE["blue"], f"{sid}_WateringCan_Body", line=PALETTE["deep_blue"], line_width=1.5)
    handle = add_shape(slide, MSO_SHAPE.ARC, x + 0.22 * scale, y - 0.27 * scale, 0.8 * scale, 0.65 * scale, PALETTE["cream"], f"{sid}_WateringCan_Handle", line=PALETTE["deep_blue"], line_width=3)
    handle.fill.background()
    spout = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + 1.15 * scale),
        Inches(y + 0.25 * scale),
        Inches(x + 2.0 * scale),
        Inches(y - 0.05 * scale),
    )
    spout.line.color.rgb = rgb(PALETTE["deep_blue"])
    spout.line.width = Pt(9 * scale)
    set_shape_meta(spout, f"{sid}_WateringCan_Spout")
    for i in range(3):
        drop = add_shape(slide, MSO_SHAPE.OVAL, x + 1.9 * scale + i * 0.16 * scale, y + 0.28 * scale + i * 0.13 * scale, 0.10 * scale, 0.18 * scale, "72C8F0", f"{sid}_WaterDrop_{i+1}")
        drop.rotation = 15


def add_sun(slide, x: float, y: float, scale: float, sid: str) -> None:
    sun = add_shape(slide, MSO_SHAPE.SUN, x, y, 1.2 * scale, 1.2 * scale, PALETTE["yellow"], f"{sid}_Sun", line="F2B944")
    add_text(slide, "☀", x + 0.24 * scale, y + 0.22 * scale, 0.72 * scale, 0.72 * scale, 24 * scale, f"{sid}_Sun_Face", color="A66D00")


def add_cloud(slide, x: float, y: float, scale: float, sid: str) -> None:
    add_shape(slide, MSO_SHAPE.CLOUD, x, y, 1.7 * scale, 0.95 * scale, "D8EEFA", f"{sid}_Cloud", line="83C6E8")
    for i in range(3):
        add_shape(slide, MSO_SHAPE.OVAL, x + (0.35 + i * 0.42) * scale, y + 0.88 * scale, 0.12 * scale, 0.22 * scale, "72C8F0", f"{sid}_Rain_{i+1}")


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, sid: str, name: str) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(PALETTE["deep_pink"])
    connector.line.width = Pt(5)
    connector.line.end_arrowhead = True
    set_shape_meta(connector, f"{sid}_{name}")


def add_picture_contain(slide, image_path: Path, x: float, y: float, w: float, h: float, name: str, descr: str) -> None:
    with Image.open(image_path) as image:
        iw, ih = image.size
    box_ratio = w / h
    image_ratio = iw / ih
    if image_ratio > box_ratio:
        height = w / image_ratio
        px, py, pw, ph = x, y + (h - height) / 2, w, height
    else:
        width = h * image_ratio
        px, py, pw, ph = x + (w - width) / 2, y, width, h
    picture = slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), Inches(pw), Inches(ph))
    set_shape_meta(picture, name, descr)


def resolve_assets(manifest_path: Path | None) -> dict[str, Path]:
    if manifest_path is None:
        return {}
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    root = manifest_path.parent
    output: dict[str, Path] = {}
    for asset in manifest.get("assets", []):
        value = asset.get("optimized_file") or asset.get("file")
        if not value:
            continue
        path = Path(value)
        if not path.is_absolute():
            path = root / path
        if path.exists():
            output[asset.get("asset_id")] = path
    return output


def visible_blocks(slide_spec: dict, variant: str) -> list[dict]:
    output = []
    for block in slide_spec.get("screen", {}).get("text_blocks", []):
        visible_in = block.get("visible_in", ["stage", "rehearsal"])
        if variant in visible_in:
            output.append(block)
    return output


def add_bilingual_lines(slide, slide_spec: dict, sid: str, theme: dict, y: float = 5.2, variant: str = "stage", x: float = 1.15, w: float = 11.0) -> None:
    blocks = visible_blocks(slide_spec, variant)
    if not blocks:
        return
    heights = 0.75 if len(blocks) == 1 else 0.55
    start_y = y if len(blocks) == 1 else y - 0.15
    for index, block in enumerate(blocks[:2]):
        lang = block.get("lang", "en")
        add_styled_text(
            slide,
            block,
            x,
            start_y + index * 0.62,
            w,
            heights,
            block.get("id") or f"{sid}_Line_{index+1}",
            theme,
            28 if lang == "en" else 24,
        )




def _zone(plan: dict, key: str, fallback: tuple[float, float, float, float]) -> tuple[float, float, float, float]:
    value = (plan.get("zones") or {}).get(key)
    if isinstance(value, list) and len(value) == 4:
        return tuple(float(v) for v in value)
    return fallback


def _add_storybook_container(slide, style: str, x: float, y: float, w: float, h: float, sid: str):
    style = style or "watercolor_blob"
    if style == "cloud_bubble":
        return add_shape(slide, MSO_SHAPE.CLOUD, x, y, w, h, "FFFDF8", f"{sid}_Storybook_Cloud", line="F0B6CE", line_width=1.6)
    if style == "speech_bubble":
        return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, x, y, w, h, "FFFDF8", f"{sid}_Storybook_SpeechBubble", line="A9C995", line_width=1.6)
    if style == "leaf_label":
        return add_shape(slide, MSO_SHAPE.HEXAGON, x, y, w, h, "E4F0D8", f"{sid}_Storybook_Leaf", line="88B96B", line_width=1.3)
    if style == "petal_cluster":
        return add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, "FCE4EE", f"{sid}_Storybook_Petal", line="F6A6C8", line_width=1.2)
    if style == "scalloped_panel":
        panel = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "FFFDF8", f"{sid}_Storybook_ScallopedPanel", line="F6A6C8", line_width=1.5)
        count = max(5, int(w / 0.55))
        for i in range(count):
            cx = x + 0.12 + i * max(0.35, (w - 0.24) / max(1, count - 1))
            add_shape(slide, MSO_SHAPE.OVAL, cx, y + h - 0.12, 0.28, 0.28, "FFFDF8", f"{sid}_Scallop_{i+1}", line="F6A6C8", line_width=0.8)
        return panel
    if style == "torn_paper":
        return add_shape(slide, MSO_SHAPE.FOLDED_CORNER, x, y, w, h, "FFFDF8", f"{sid}_Storybook_TornPaper", line="BFD7B7", line_width=1.4)
    # Watercolor blob fallback: a large soft oval with transparent edge.
    shape = add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, "FFF7E6", f"{sid}_Storybook_WatercolorBlob", line="D8E8C9", line_width=1.0)
    shape.fill.transparency = 8
    return shape


def _render_keyword_group(slide, blocks: list[dict], plan: dict, sid: str, theme: dict) -> None:
    if not blocks:
        return
    x, y, w, h = _zone(plan, "accent", (6.6, 4.8, 5.4, 1.5))
    count = len(blocks)
    chip_w = min(1.65, max(1.05, (w - 0.18 * (count - 1)) / max(1, count)))
    gap = 0.18
    total = chip_w * count + gap * (count - 1)
    start = x + max(0, (w - total) / 2)
    styles = [MSO_SHAPE.OVAL, MSO_SHAPE.HEXAGON, MSO_SHAPE.OVAL]
    fills = ["FCE1EC", "E5F2D8", "FFF0B8"]
    for i, block in enumerate(blocks):
        bx = start + i * (chip_w + gap)
        add_shape(slide, styles[i % len(styles)], bx, y, chip_w, min(h, 0.82), fills[i % len(fills)], f"{sid}_KeywordChip_{i+1}", line="E3B7C8", line_width=1.0)
        add_styled_text(slide, block, bx + 0.08, y + 0.08, chip_w - 0.16, min(h, 0.66), block.get("id") or f"{sid}_Keyword_{i+1}", theme, 22)


def render_storybook_plan(slide, spec: dict, sid: str, assets: dict[str, Path], variant: str, theme: dict) -> bool:
    visual = spec.get("visual", {}) or {}
    plan = visual.get("layout_plan") or {}
    if not plan:
        return False
    archetype = plan.get("archetype", "picturebook-window")
    asset_ids = visual.get("asset_ids", []) or []
    image_path = next((assets.get(asset_id) for asset_id in asset_ids if assets.get(asset_id)), None)
    blocks = visible_blocks(spec, variant)
    keyword_blocks = [b for b in blocks if b.get("role") in {"keyword", "label", "vocabulary"}]
    main_blocks = [b for b in blocks if b not in keyword_blocks]

    if archetype == "weather-triptych":
        lx, ly, lw, lh = _zone(plan, "left", (0.45, 0.75, 3.65, 4.55))
        cx, cy, cw, ch = _zone(plan, "center", (4.55, 0.55, 4.2, 5.1))
        rx, ry, rw, rh = _zone(plan, "right", (9.15, 0.75, 3.65, 4.55))
        _add_storybook_container(slide, "watercolor_blob", lx, ly, lw, lh, sid + "_L")
        _add_storybook_container(slide, "watercolor_blob", cx, cy, cw, ch, sid + "_C")
        _add_storybook_container(slide, "watercolor_blob", rx, ry, rw, rh, sid + "_R")
        add_sun(slide, lx + 0.8, ly + 0.7, 1.4, sid)
        add_flower(slide, cx + 1.5, cy + 0.75, 1.5, sid, "WeatherFlower", bloom=False)
        add_cloud(slide, rx + 0.7, ry + 0.8, 1.35, sid)
    elif archetype == "three-vignette-growth":
        for n, key in enumerate(("v1", "v2", "v3"), 1):
            x, y, w, h = _zone(plan, key, (0.6 + (n-1)*4.2, 0.8, 3.5, 4.6))
            _add_storybook_container(slide, "scalloped_panel", x, y, w, h, f"{sid}_V{n}")
            add_flower(slide, x + 0.9, y + 0.65, 1.25 + 0.12*n, sid, f"Growth{n}", bloom=n == 3)
    else:
        hx, hy, hw, hh = _zone(plan, "hero", (0.55, 0.55, 6.4, 6.1))
        if image_path:
            add_picture_contain(slide, image_path, hx, hy, hw, hh, f"{sid}_Storybook_Hero", "Picture-book hero or story object")
        else:
            _add_storybook_container(slide, "watercolor_blob", hx, hy, hw, hh, sid + "_HeroBackdrop")
            add_flower(slide, hx + hw * 0.34, hy + hh * 0.18, min(2.0, hw / 2.6), sid, "StorybookHeroFlower", bloom=archetype not in {"diagonal-action"})
            if archetype == "diagonal-action":
                add_watering_can(slide, hx + 0.25, hy + hh * 0.43, 1.0, sid)

    tx, ty, tw, th = _zone(plan, "text", (7.0, 0.9, 5.5, 4.2))
    _add_storybook_container(slide, plan.get("text_container_style", "watercolor_blob"), tx, ty, tw, th, sid)
    if main_blocks:
        block_count = max(1, min(3, len(main_blocks)))
        line_h = min(1.15, max(0.48, (th - 0.30) / block_count))
        if line_h * block_count > th - 0.16:
            line_h = max(0.42, (th - 0.16) / block_count)
        start_y = ty + max(0.08, (th - line_h * block_count) / 2)
        text_h = max(0.38, line_h - 0.05)
        for i, block in enumerate(main_blocks[:3]):
            add_styled_text(slide, block, tx + 0.28, start_y + i * line_h, tw - 0.56, text_h, block.get("id") or f"{sid}_PlannedText_{i+1}", theme, 30 if block.get("lang") == "en" else 25)
    _render_keyword_group(slide, keyword_blocks, plan, sid, theme)
    return True

def render_layout(slide, spec: dict, sid: str, assets: dict[str, Path], variant: str, theme: dict) -> None:
    if render_storybook_plan(slide, spec, sid, assets, variant, theme):
        return
    visual = spec.get("visual", {})
    layout = visual.get("layout", "show-one-object")
    asset_ids = visual.get("asset_ids", [])
    image_path = next((assets.get(asset_id) for asset_id in asset_ids if assets.get(asset_id)), None)

    if layout == "kindergarten-speech-cover":
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 0.75, 5.6, 4.4, "FFFDF8", f"{sid}_Title_Panel", line="F6D4E2", line_width=1.5)
        blocks = visible_blocks(spec, variant)
        for index, block in enumerate(blocks[:2]):
            add_styled_text(
                slide,
                block,
                1.2,
                1.45 + index * 1.25,
                4.9,
                1.0,
                block.get("id") or f"{sid}_Title_{index+1}",
                theme,
                38 if index == 0 else 28,
            )
        if image_path:
            add_picture_contain(slide, image_path, 6.65, 0.45, 6.1, 6.35, f"{sid}_Hero_Image", "Consistent child speaker avatar in a flower garden")
        else:
            add_flower(slide, 8.0, 1.5, 1.8, sid, "CoverFlower", bloom=True)
        return

    if layout in {"hello-and-wave", "this-is-me"}:
        if image_path:
            add_picture_contain(slide, image_path, 0.7, 0.55, 6.1, 5.75, f"{sid}_Speaker_Avatar", "Child speaker avatar")
        else:
            add_shape(slide, MSO_SHAPE.OVAL, 1.2, 1.05, 3.4, 3.4, "FFE2EE", f"{sid}_Avatar_Placeholder")
            add_text(slide, "👋", 2.0, 1.75, 1.8, 1.8, 60, f"{sid}_Wave_Icon")
        add_shape(slide, MSO_SHAPE.CLOUD, 7.0, 1.1, 5.2, 2.8, "FFFDF8", f"{sid}_Speech_Cloud", line="F6D4E2", line_width=1.5)
        add_bilingual_lines(slide, spec, sid, theme, y=2.0, variant=variant, x=7.35, w=4.45)
        return

    if layout == "show-one-object":
        add_shape(slide, MSO_SHAPE.OVAL, 0.8, 0.75, 5.4, 4.8, "FFF1D7", f"{sid}_Object_Backdrop")
        if image_path:
            add_picture_contain(slide, image_path, 1.0, 0.9, 5.0, 4.4, f"{sid}_Hero_Object", "Story object")
        else:
            add_flower(slide, 2.05, 1.25, 1.8, sid, "HeroFlower", bloom=True)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.7, 1.1, 5.6, 3.7, "FFFDF8", f"{sid}_Story_Card", line="CFE8D0", line_width=1.5)
        add_bilingual_lines(slide, spec, sid, theme, y=2.0, variant=variant, x=7.1, w=4.8)
        return

    if layout == "one-action-one-line":
        add_watering_can(slide, 1.1, 1.35, 1.15, sid)
        add_flower(slide, 4.15, 1.05, 1.5, sid, "WateredFlower", bloom=False)
        add_shape(slide, MSO_SHAPE.HEART, 8.85, 1.25, 1.1, 1.0, "FFD6E5", f"{sid}_Action_Heart")
        add_bilingual_lines(slide, spec, sid, theme, y=5.25, variant=variant)
        return

    if layout == "weather-pair":
        add_sun(slide, 1.0, 1.0, 1.45, sid)
        add_cloud(slide, 9.8, 1.05, 1.25, sid)
        add_flower(slide, 5.65, 1.4, 1.45, sid, "WeatherFlower", bloom=False)
        add_styled_text(slide, {"text": "sunshine", "lang": "en", "role": "label", "font_size_pt": 25}, 0.75, 3.05, 2.3, 0.55, f"{sid}_Sun_Label", theme, 25)
        add_styled_text(slide, {"text": "rain", "lang": "en", "role": "label", "font_size_pt": 25}, 10.0, 3.05, 1.8, 0.55, f"{sid}_Rain_Label", theme, 25)
        add_bilingual_lines(slide, spec, sid, theme, y=5.25, variant=variant)
        return

    if layout == "before-and-after":
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 0.75, 4.7, 4.45, "FFFDF8", f"{sid}_Before_Card", line="F2D6E1")
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.75, 0.75, 4.7, 4.45, "FFFDF8", f"{sid}_After_Card", line="D5EACF")
        add_flower(slide, 2.0, 1.35, 1.5, sid, "BeforeFlower", bloom=False)
        add_flower(slide, 8.85, 1.05, 1.7, sid, "AfterFlower", bloom=True)
        add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 5.75, 2.47, 1.7, 0.78, "F6A6C8", f"{sid}_Change_Arrow", line="D95D8A", line_width=1.2)
        # Labels live in the clean title band of each card, not on top of the flower/pot.
        add_styled_text(slide, {"text": "before", "lang": "en", "role": "label", "font_size_pt": 23}, 1.35, 0.28, 2.2, 0.5, f"{sid}_Before_Label", theme, 23)
        add_styled_text(slide, {"text": "after", "lang": "en", "role": "label", "font_size_pt": 23}, 8.42, 0.28, 2.2, 0.5, f"{sid}_After_Label", theme, 23)
        add_bilingual_lines(slide, spec, sid, theme, y=5.5, variant=variant)
        return

    if layout == "emotion-moment":
        for i, (x, y, scale) in enumerate([(1.0, 1.0, 1.4), (10.8, 1.1, 0.9), (9.4, 4.0, 0.65)]):
            add_shape(slide, MSO_SHAPE.HEART, x, y, 1.1 * scale, 1.0 * scale, "FFD0E2", f"{sid}_Heart_{i+1}")
        add_flower(slide, 5.25, 0.75, 2.0, sid, "HappyFlower", bloom=True)
        add_bilingual_lines(slide, spec, sid, theme, y=5.35, variant=variant)
        return

    if layout == "simple-closing":
        if image_path:
            add_picture_contain(slide, image_path, 6.6, 0.5, 6.0, 6.2, f"{sid}_Closing_Avatar", "Consistent child speaker avatar waving goodbye")
        else:
            add_flower(slide, 8.2, 1.3, 1.9, sid, "ClosingFlower", bloom=True)
        add_shape(slide, MSO_SHAPE.CLOUD, 0.85, 1.15, 5.7, 3.5, "FFFDF8", f"{sid}_Closing_Cloud", line="F6D4E2")
        add_bilingual_lines(slide, spec, sid, theme, y=2.05, variant=variant, x=1.35, w=4.8)
        return

    add_flower(slide, 1.45, 1.15, 1.7, sid, "FallbackFlower", bloom=True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.25, 1.0, 6.8, 3.8, "FFFDF8", f"{sid}_Fallback_Card", line="D5EACF")
    add_bilingual_lines(slide, spec, sid, theme, y=2.0, variant=variant)


def add_rehearsal_overlay(slide, spec: dict, sid: str) -> None:
    speaker = spec.get("speaker", {})
    actions = ", ".join(speaker.get("actions", []))
    cue = " / ".join(filter(None, [speaker.get("script_en", ""), speaker.get("script_zh", "")]))
    if actions:
        cue += f"\n动作：{actions}"
    panel = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.45, 6.25, 12.43, 0.92, "503F45", f"{sid}_Rehearsal_Panel")
    panel.fill.transparency = 5
    add_text(slide, cue, 0.72, 6.37, 11.9, 0.62, 15, f"{sid}_Rehearsal_Cue", lang="zh", color=PALETTE["white"], bold=False, align=PP_ALIGN.LEFT)


def build_deck(spec: dict, output: Path, assets: dict[str, Path], variant: str) -> None:
    theme = load_typography_theme(spec)
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    for slide_spec in spec.get("slides", []):
        sid = slide_spec.get("id", "S")
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        add_soft_background(slide, sid)
        add_visual_components(slide, slide_spec, sid, layer="background")
        render_layout(slide, slide_spec, sid, assets, variant, theme)
        add_visual_components(slide, slide_spec, sid, layer="foreground")
        if variant == "rehearsal":
            add_rehearsal_overlay(slide, slide_spec, sid)
        speaker = slide_spec.get("speaker", {})
        notes_parts = []
        if speaker.get("script_en"):
            notes_parts.append("EN: " + speaker["script_en"])
        if speaker.get("script_zh"):
            notes_parts.append("ZH: " + speaker["script_zh"])
        if speaker.get("actions"):
            notes_parts.append("ACTIONS: " + "; ".join(speaker["actions"]))
        if speaker.get("pronunciation"):
            notes_parts.append("PRONUNCIATION: " + str(speaker["pronunciation"]))
        slide.notes_slide.notes_text_frame.text = "\n".join(notes_parts)
    presentation.core_properties.title = spec.get("title", "Child Speech")
    presentation.core_properties.subject = "Child stage speech presentation"
    presentation.core_properties.author = "kv-kids-learning-slides"
    presentation.save(output)
    if any(slide.get("stage", {}).get("animation_mode") not in {None, "none"} for slide in spec.get("slides", [])):
        from apply_gentle_transitions import apply_transitions
        apply_transitions(output)


def count_english_words(text: str) -> int:
    import re
    return len(re.findall(r"\b[A-Za-z]+(?:['’][A-Za-z]+)?\b", text))


def first_screen_line(slide: dict, lang: str) -> str:
    for block in slide.get("screen", {}).get("text_blocks", []):
        if block.get("lang") == lang and block.get("text"):
            return block["text"].splitlines()[0].strip()
    return ""


def write_script(spec: dict, output: Path) -> None:
    profile = spec.get("speech_profile", {})
    audience = spec.get("audience", {})
    speaker_name = spec.get("speaker_name") or spec.get("child_name") or "________"
    age_value = spec.get("speaker_age") or audience.get("age") or audience.get("age_band") or "________"
    total_words = sum(count_english_words(slide.get("speaker", {}).get("script_en", "")) for slide in spec.get("slides", []))
    visible_words = spec.get("stage_visible_english_words")
    if visible_words is None:
        visible_words = sum(
            count_english_words(block.get("text", ""))
            for slide in spec.get("slides", [])
            for block in slide.get("screen", {}).get("text_blocks", [])
            if block.get("lang") == "en" and "stage" in block.get("visible_in", ["stage", "rehearsal"])
        )
    duration_seconds = profile.get("target_duration_seconds") or spec.get("speech_duration_seconds")
    if not duration_seconds:
        duration_seconds = sum(int(slide.get("stage", {}).get("duration_seconds", 0) or 0) for slide in spec.get("slides", []))
    duration_text = f"约 {max(1, round(duration_seconds / 60, 1))} 分钟" if duration_seconds else "________"
    bilingual_mode = spec.get("bilingual", {}).get("mode", "________")
    visual_system = spec.get("visual_system", "________")
    title = spec.get("title", "Speaker Script")
    lines = [
        f"# {title}",
        "",
        f"> **Speaker / 演讲者：** {speaker_name}",
        f"> **Age / 年龄：** {age_value}",
        f"> **Duration / 建议时长：** {duration_text}",
        f"> **Difficulty / 难度：** {profile.get('difficulty', '________')}",
        f"> **Bilingual mode / 双语模式：** {bilingual_mode}",
        f"> **English words / 英文词数：** {total_words}",
        f"> **Stage cue words / PPT 可见英文提示词：** {visible_words}",
        f"> **Visual style / 视觉风格：** {visual_system}",
        "",
        "---",
        "",
        "## 使用建议 / How to Use",
        "- 先按 PPT 页码逐段排练，再使用文末连贯稿完整练习。",
        "- 英文是推荐表达，中文用于理解和家长辅助，可按比赛规则调整语言顺序。",
        "- 动作与节奏提示只用于表演训练，不需要逐字念出。",
        "- 建议先修改这份 Markdown，再重新生成 Word，保证两个版本一致。",
        "",
    ]
    continuous_en = []
    continuous_zh = []
    for slide in spec.get("slides", []):
        speaker = slide.get("speaker", {})
        en = speaker.get("script_en", "").strip()
        zh = speaker.get("script_zh", "").strip()
        actions = speaker.get("actions", []) or []
        pronunciation = speaker.get("pronunciation", []) or slide.get("speaker", {}).get("rhythm", []) or []
        if isinstance(pronunciation, str):
            pronunciation = [pronunciation]
        en_title = first_screen_line(slide, "en")
        zh_title = first_screen_line(slide, "zh")
        display_title = " / ".join(filter(None, [en_title, zh_title])) or slide.get("role", "Speech beat")
        lines.extend([
            f"## {slide.get('id')} · {display_title}",
            "",
            "### English",
            en or "[Edit English script here]",
            "",
            "### 中文",
            zh or "[在这里修改中文参考稿]",
            "",
            "### Stage actions / 舞台动作",
        ])
        if actions:
            lines.extend(f"- {action}" for action in actions)
        else:
            lines.append("- No special action / 无特别动作")
        lines.extend(["", "### Rhythm and pronunciation / 节奏与发音"])
        if pronunciation:
            lines.extend(f"- {item}" for item in pronunciation)
        else:
            duration = slide.get("stage", {}).get("duration_seconds")
            cue = f"Speak slowly; allow about {duration} seconds for this page." if duration else "Speak slowly and pause naturally."
            lines.append(f"- {cue}")
        lines.extend([
            "",
            "### Editable notes / 可修改备注",
            "> 家长、老师或孩子可以在这里补充个性化提醒。",
            "",
        ])
        if en:
            continuous_en.append(en)
        if zh:
            continuous_zh.append(zh)
    lines.extend([
        "---",
        "",
        "## Continuous English Script / 英文连贯稿",
        "",
        "\n\n".join(continuous_en),
        "",
        "## 中文参考稿",
        "",
        "\n\n".join(continuous_zh),
        "",
    ])
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")

def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("spec")
    parser.add_argument("--stage-out", required=True)
    parser.add_argument("--rehearsal-out")
    parser.add_argument("--script-out")
    parser.add_argument("--script-docx-out")
    parser.add_argument("--script-theme", choices=["paper-theater", "watercolor-glow", "clean-editorial"], default="paper-theater")
    parser.add_argument("--assets")
    args = parser.parse_args()

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    assets = resolve_assets(Path(args.assets) if args.assets else None)
    build_deck(spec, Path(args.stage_out), assets, "stage")
    if args.rehearsal_out:
        build_deck(spec, Path(args.rehearsal_out), assets, "rehearsal")
    script_path = Path(args.script_out) if args.script_out else None
    if args.script_docx_out and script_path is None:
        script_path = Path(args.script_docx_out).with_suffix(".md")
    if script_path:
        write_script(spec, script_path)
    if args.script_docx_out:
        from build_speaker_script_docx import build_docx_from_markdown
        build_docx_from_markdown(script_path, Path(args.script_docx_out), args.script_theme)
    print(
        json.dumps(
            {
                "stage": args.stage_out,
                "rehearsal": args.rehearsal_out,
                "script": str(script_path) if script_path else None,
                "script_docx": args.script_docx_out,
                "slides": len(spec.get("slides", [])),
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
