#!/usr/bin/env python3
"""Portable fallback: build a simple editable PPTX from slide_spec.json."""
import argparse, json, sys
from pathlib import Path
try:
 from pptx import Presentation
 from pptx.util import Inches, Pt
 from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
 from pptx.enum.shapes import MSO_SHAPE
 from pptx.dml.color import RGBColor
except Exception as e:
 print('python-pptx is required: pip install python-pptx', file=sys.stderr); raise

def rgb(h): h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
def add_text(slide, text, x,y,w,h,size=24,bold=False,color='#273043',align=PP_ALIGN.LEFT):
 box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.name='editable-text'
 tf=box.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
 p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=str(text); r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(color); r.font.name='Microsoft YaHei'
 return box

def build(spec, out):
 prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
 for idx,s in enumerate(spec['slides'],1):
  slide=prs.slides.add_slide(prs.slide_layouts[6]); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb('#FFF9ED')
  # top accent
  sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.16)); sh.fill.solid(); sh.fill.fore_color.rgb=rgb('#FFB703'); sh.line.fill.background(); sh.name=f"{s.get('id','S')}-accent"
  title=s.get('screen',{}).get('title',''); add_text(slide,title,.7,.45,11.9,.75,30,True)
  body=s.get('screen',{}).get('body') or s.get('screen',{}).get('prompt') or ''
  # Native card
  card=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.7),Inches(1.45),Inches(11.9),Inches(4.65)); card.fill.solid(); card.fill.fore_color.rgb=rgb('#FFFFFF'); card.line.color.rgb=rgb('#DDE7EA'); card.name=f"{s.get('id','S')}-content-card"
  add_text(slide,body,1.05,1.75,11.2,3.45,24,False)
  goal=s.get('learning_goal',''); add_text(slide,'学习任务：'+goal,1.05,5.38,10.6,.5,15,False,'#2A6F97')
  add_text(slide,f'{idx}/{len(spec["slides"])}',11.75,6.75,.7,.35,12,False,'#6B7280',PP_ALIGN.RIGHT)
 prs.save(out)

def main():
 p=argparse.ArgumentParser(); p.add_argument('spec'); p.add_argument('output'); a=p.parse_args(); spec=json.loads(Path(a.spec).read_text(encoding='utf-8')); build(spec,a.output); print(a.output)
if __name__=='__main__': main()
